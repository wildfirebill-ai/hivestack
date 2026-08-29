"""Office document engine — Word (python-docx), Excel (openpyxl), PowerPoint
(python-pptx). Build → registry → preview → audit → diff. Fully offline."""

from __future__ import annotations

import difflib
import re
import uuid
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from ..config import settings
from ..db import _conn

_EXT = {"word": "docx", "spreadsheet": "xlsx", "slides": "pptx"}


def docs_dir() -> Path:
    p = settings.data_dir / "docs"
    p.mkdir(parents=True, exist_ok=True)
    return p


def _register(name: str, fmt: str, relpath: str) -> str:
    doc_id = uuid.uuid4().hex[:12]
    with _conn() as con:
        con.execute("INSERT INTO documents(id, name, format, path) VALUES (?,?,?,?)", (doc_id, name, fmt, relpath))
    return doc_id


def _path(rel: str) -> Path:
    return settings.data_dir / rel


def list_docs() -> list[dict]:
    with _conn() as con:
        rows = con.execute("SELECT id, name, format, path, created_at FROM documents ORDER BY created_at DESC").fetchall()
    return [dict(r) for r in rows]


def get_doc(doc_id: str) -> dict | None:
    with _conn() as con:
        row = con.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
    return dict(row) if row else None


def _sub(text: str, data: dict | None) -> str:
    if not data:
        return text
    for k, v in data.items():
        text = text.replace("{{" + k + "}}", str(v))
    return text


# ------------------------------------------------------------------ build
def build_word(
    title: str,
    sections: list[dict],
    table_rows: list[list] | None = None,
    merge_data: dict | None = None,
) -> dict:
    doc = Document()
    doc.add_heading(_sub(title, merge_data), 0)
    for sec in sections or []:
        doc.add_heading(_sub(str(sec.get("heading", "")), merge_data), level=1)
        body = sec.get("body", "")
        if isinstance(body, list):
            for item in body:
                doc.add_paragraph(_sub(str(item), merge_data), style="List Bullet")
        elif body:
            doc.add_paragraph(_sub(str(body), merge_data))
    if table_rows:
        ncols = max(len(r) for r in table_rows)
        table = doc.add_table(rows=len(table_rows), cols=ncols)
        table.style = "Light Grid Accent 1"
        for i, row in enumerate(table_rows):
            for j in range(ncols):
                table.rows[i].cells[j].text = _sub(str(row[j]) if j < len(row) else "", merge_data)
    path = docs_dir() / f"{uuid.uuid4().hex[:8]}.docx"
    doc.save(path)
    doc_id = _register(title, "word", str(path.relative_to(settings.data_dir)))
    return {"id": doc_id, "format": "word", "path": str(path.relative_to(settings.data_dir)), "audit": audit(path)}


def build_spreadsheet(sheets: list[dict]) -> dict:
    wb = Workbook()
    wb.remove(wb.active)
    for sheet in sheets:
        ws = wb.create_sheet(title=str(sheet.get("name", "Sheet"))[:31])
        rows = sheet.get("rows", [])
        width = sheet.get("width", 16)
        for r in rows:
            ws.append(r)
        for col in ws.columns:
            letter = col[0].column_letter
            ws.column_dimensions[letter].width = width
    path = docs_dir() / f"{uuid.uuid4().hex[:8]}.xlsx"
    wb.save(path)
    doc_id = _register("spreadsheet", "spreadsheet", str(path.relative_to(settings.data_dir)))
    return {"id": doc_id, "format": "spreadsheet", "path": str(path.relative_to(settings.data_dir)), "audit": audit(path)}


def build_slides(title: str, slides: list[dict]) -> dict:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = title
    bullet_layout = prs.slide_layouts[1]
    for s in slides or []:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = str(s.get("title", ""))
        body = slide.placeholders[1].text_frame
        for i, bullet in enumerate(s.get("bullets", []) or []):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(bullet)
    path = docs_dir() / f"{uuid.uuid4().hex[:8]}.pptx"
    prs.save(path)
    doc_id = _register(title, "slides", str(path.relative_to(settings.data_dir)))
    return {"id": doc_id, "format": "slides", "path": str(path.relative_to(settings.data_dir)), "audit": audit(path)}


# ------------------------------------------------------------------ preview / audit / diff
def preview(path: Path) -> dict:
    fmt = path.suffix.lower()
    lines: list[str] = []
    if fmt == ".docx":
        d = Document(path)
        for p in d.paragraphs:
            if p.text.strip():
                lines.append(p.text)
        for t in d.tables:
            if t.rows:
                lines.append("| " + " | ".join(c.text for c in t.rows[0].cells) + " |")
                for row in t.rows[1:]:
                    lines.append("| " + " | ".join(c.text for c in row.cells) + " |")
        return {"format": "word", "lines": lines}
    if fmt == ".xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(path, data_only=False)
        for ws in wb.worksheets:
            lines.append(f"# sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v) for v in row]
                if any(vals):
                    lines.append(" | ".join(vals))
        return {"format": "spreadsheet", "lines": lines}
    if fmt == ".pptx":
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text.strip()
                    if txt:
                        lines.append(txt)
        return {"format": "slides", "lines": lines}
    return {"format": fmt, "lines": []}


def audit(path: Path) -> dict:
    pre = preview(path)
    checks = [
        {"name": "producible", "pass": len(pre["lines"]) > 0, "detail": f"{len(pre['lines'])} content lines"},
        {"name": "heading_present", "pass": any(l.strip() for l in pre["lines"][:3]), "detail": "top content"},
        {"name": "size_ok", "pass": path.stat().st_size > 0, "detail": f"{path.stat().st_size} bytes"},
    ]
    score = round(sum(1 for c in checks if c["pass"]) / len(checks), 2)
    return {"checks": checks, "score": score, "pass": score == 1.0}


def diff(a_path: Path, b_path: Path) -> dict:
    sa = preview(a_path)["lines"]
    sb = preview(b_path)["lines"]
    delta = list(difflib.unified_diff(sa, sb, lineterm="", n=1))
    return {"added": sum(1 for l in delta if l.startswith("+") and not l.startswith("+++")),
            "removed": sum(1 for l in delta if l.startswith("-") and not l.startswith("---")),
            "diff": delta[:120]}


def mail_merge_tokens(text: str) -> list[str]:
    return re.findall(r"\{\{(\w+)\}\}", text)